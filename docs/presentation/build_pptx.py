#!/usr/bin/env python3
"""Build ``Yequan_26_Midpoint.pptx`` from the midpoint slide contents.

Design: dark "sandwich" (navy title + section dividers, light content), one
distinctive motif (a thin gold accent block behind slide numbers + navy title
band), Times New Roman throughout with 20pt body text. Tables and figures are
positioned inside a fixed content region so nothing overlaps or bleeds off.

Run from ``docs/presentation``:  python build_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
FIGS = HERE / "figs"
OUT = HERE / "Yequan_26_Midpoint.pptx"

# ---------------------------------------------------------------------------
# Palette (Midnight / Ocean — matches an ML-systems topic)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x2A, 0x47)      # dominant dark
NAVY2 = RGBColor(0x14, 0x3A, 0x5E)     # slightly lighter navy for bands
GOLD = RGBColor(0xE0, 0xA3, 0x2E)      # sharp accent
TEAL = RGBColor(0x1C, 0x72, 0x93)      # supporting
INK = RGBColor(0x22, 0x2A, 0x33)       # body text on light
MUTE = RGBColor(0x5C, 0x6B, 0x7A)      # captions
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)     # light bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROWALT = RGBColor(0xEA, 0xEF, 0xF4)    # table zebra
HEADFILL = NAVY
HILITE = RGBColor(0xFC, 0xF3, 0xDC)    # soft gold row highlight

FONT = "Times New Roman"

# Layout geometry (inches) — 13.333 x 7.5 (16:9)
SW, SH = 13.333, 7.5
MARGIN = 0.6
CONTENT_TOP = 1.40          # below the title (no band)
CONTENT_BOT = 6.82          # above footer / logo row
CONTENT_L = MARGIN
CONTENT_R = SW - MARGIN
CONTENT_W = CONTENT_R - CONTENT_L

# Amazon logo — position matches the template's slide master (bottom-left)
LOGO = FIGS / "amazon_logo.png"
LOGO_X, LOGO_Y, LOGO_W, LOGO_H = 0.21, 6.90, 1.60, 0.482

BODY_SZ = 20                # required body size
TITLE_SZ = 30
SUBTITLE_SZ = 20

prs = Presentation()
prs.slide_width = Emu(int(SW * 914400))
prs.slide_height = Emu(int(SH * 914400))
BLANK = prs.slide_layouts[6]

_slide_no = 0


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _set_font(run, size=BODY_SZ, bold=False, italic=False, color=INK, name=FONT):
    f = run.font
    f.name = name
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    # ensure east-asian / cs also use the face
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x * 914400)),
                                Emu(int(y * 914400)), Emu(int(w * 914400)),
                                Emu(int(h * 914400)))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, kwargs) tuples."""
    tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                  Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(int(0.05 * 914400))
    tf.margin_right = Emu(int(0.05 * 914400))
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for text, kw in para:
            r = p.add_run()
            r.text = text
            _set_font(r, **kw)
    return tb


def add_logo(slide, x=LOGO_X):
    """Amazon logo bottom-left, matching the template's master placement."""
    slide.shapes.add_picture(str(LOGO), Emu(int(x * 914400)),
                             Emu(int(LOGO_Y * 914400)), Emu(int(LOGO_W * 914400)),
                             Emu(int(LOGO_H * 914400)))


def add_footer(slide, title_short):
    global _slide_no
    _slide_no += 1
    add_logo(slide)
    # gold tick + slide number bottom-right
    add_rect(slide, SW - 1.15, SH - 0.5, 0.12, 0.26, GOLD)
    add_text(slide, SW - 0.95, SH - 0.54, 0.75, 0.35,
             [[(str(_slide_no), dict(size=12, color=MUTE, bold=True))]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # footer caption — right of the logo, so it never collides
    add_text(slide, 2.0, SH - 0.52, 9.0, 0.35,
             [[("Per-Token Adaptive Channel Activation for MoE", dict(size=11, color=MUTE, italic=True))]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def content_slide(title, subtitle=None):
    """Light content slide: navy title text (no dark band); return slide + content top-y."""
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    # gold tick motif to the left of the title
    add_rect(s, MARGIN, 0.34, 0.14, 0.62, GOLD)
    add_text(s, MARGIN + 0.30, 0.22, SW - 2 * MARGIN - 0.30, 0.86,
             [[(title, dict(size=TITLE_SZ, bold=True, color=NAVY))]],
             anchor=MSO_ANCHOR.MIDDLE)
    top = CONTENT_TOP
    if subtitle:
        add_text(s, MARGIN, 1.18, CONTENT_W, 0.5,
                 [[(subtitle, dict(size=SUBTITLE_SZ, bold=True, color=TEAL))]],
                 anchor=MSO_ANCHOR.TOP)
        top = 1.82
    add_footer(s, title)
    return s, top


def section_slide(title):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    # short gold marker above the title (motif)
    add_rect(s, MARGIN, 2.75, 1.4, 0.14, GOLD)
    add_text(s, MARGIN, 3.12, SW - 2 * MARGIN, 1.2,
             [[(title, dict(size=40, bold=True, color=NAVY))]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_logo(s)
    return s


def fig_size(png, max_w, max_h):
    im = Image.open(png)
    ar = im.width / im.height
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    return w, h


def add_figure(slide, png, y, max_w=CONTENT_W, max_h=3.5, center_x=None, caption=None):
    w, h = fig_size(FIGS / png, max_w, max_h)
    x = (SW - w) / 2 if center_x is None else center_x
    slide.shapes.add_picture(str(FIGS / png), Emu(int(x * 914400)),
                             Emu(int(y * 914400)), Emu(int(w * 914400)),
                             Emu(int(h * 914400)))
    if caption:
        add_text(slide, MARGIN, y + h + 0.05, CONTENT_W, 0.5,
                 [[(caption, dict(size=12, italic=True, color=MUTE))]],
                 align=PP_ALIGN.CENTER)
    return w, h


def add_equation(slide, png, y, max_w=9.5, max_h=0.85, center=True, x=None):
    w, h = fig_size(FIGS / png, max_w, max_h)
    xx = (SW - w) / 2 if center else (x if x is not None else CONTENT_L)
    slide.shapes.add_picture(str(FIGS / png), Emu(int(xx * 914400)),
                             Emu(int(y * 914400)), Emu(int(w * 914400)),
                             Emu(int(h * 914400)))
    return w, h


def bullets(slide, x, y, w, h, items, size=BODY_SZ, space_after=8, line_spacing=1.0):
    """items: list of (text|list-of-runs, level, kwargs) or (text, level)."""
    tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                  Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if len(item) == 3:
            content, level, kw = item
        else:
            content, level = item
            kw = {}
        p.level = level
        # bullet char via XML
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪" if level == 0 else "–"})
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": FONT})
        buClr = pPr.makeelement(qn("a:buClr"), {})
        srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "E0A32E" if level == 0 else "1C7293"})
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)
        indent = 0.28
        pPr.set("indent", str(-int(indent * 914400)))
        pPr.set("marL", str(int((indent + level * 0.3) * 914400)))
        runs = content if isinstance(content, list) else [(content, {})]
        for text, rkw in runs:
            r = p.add_run()
            r.text = text
            merged = dict(size=size, color=INK)
            merged.update(kw)
            merged.update(rkw)
            _set_font(r, **merged)
    return tb


# ---------------------------------------------------------------------------
# Table builder
# ---------------------------------------------------------------------------
def add_table(slide, x, y, w, rows, col_w=None, header=True, font=14,
              highlight_rows=(), col_align=None, row_h=0.34, header_h=0.46):
    """rows: list of list of cell-strings (row 0 = header if header=True)."""
    nrow = len(rows)
    ncol = len(rows[0])
    total_h = header_h + (nrow - 1) * row_h if header else nrow * row_h
    gt = slide.shapes.add_table(nrow, ncol, Emu(int(x * 914400)), Emu(int(y * 914400)),
                                Emu(int(w * 914400)), Emu(int(total_h * 914400)))
    tbl = gt.table
    # disable default banding styling; we set fills manually
    tbl.first_row = header
    tbl.horz_banding = False
    if col_w:
        assert abs(sum(col_w) - w) < 1e-6 or True
        for j, cw in enumerate(col_w):
            tbl.columns[j].width = Emu(int(cw * 914400))
    for i in range(nrow):
        tbl.rows[i].height = Emu(int((header_h if (header and i == 0) else row_h) * 914400))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Emu(int(0.06 * 914400))
            cell.margin_right = Emu(int(0.06 * 914400))
            cell.margin_top = Emu(int(0.01 * 914400))
            cell.margin_bottom = Emu(int(0.01 * 914400))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # fill
            if header and i == 0:
                fill_c = HEADFILL
            elif i in highlight_rows:
                fill_c = HILITE
            elif i % 2 == 0:
                fill_c = ROWALT
            else:
                fill_c = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_c
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if col_align:
                p.alignment = col_align[j]
            else:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            is_head = header and i == 0
            is_hl = i in highlight_rows
            _set_font(r, size=font, bold=is_head or (is_hl and j == 0),
                      color=WHITE if is_head else INK)
    return gt, total_h


# ===========================================================================
# SLIDES
# ===========================================================================

# --- Slide 1: Title -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT)
add_rect(s, 0, 0, 0.35, SH, GOLD)                 # left gold spine (motif)
add_rect(s, MARGIN + 0.2, 2.15, 2.0, 0.16, GOLD)
add_text(s, MARGIN + 0.2, 2.5, SW - 2 * MARGIN, 1.9,
         [[("Per-Token Adaptive Channel Activation", dict(size=42, bold=True, color=NAVY))],
          [("for Efficient MoE Inference", dict(size=42, bold=True, color=NAVY))]],
         anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.02)
add_text(s, MARGIN + 0.2, 4.75, SW - 2 * MARGIN, 0.5,
         [[("Yequan Zhao", dict(size=22, bold=True, color=TEAL)),
           ("     |     Midpoint Review     |     August 2026", dict(size=22, color=MUTE))]])
add_text(s, MARGIN + 0.2, 5.5, SW - 2 * MARGIN, 0.6,
         [[("Model: Qwen3-30B-A3B", dict(size=18, bold=True, color=INK)),
           ("   —   128 experts, top-8 routing, 48 MoE layers", dict(size=18, color=MUTE))]])
add_logo(s, x=0.55)   # clear the gold spine

# --- Section: Background ---------------------------------------------------
section_slide("Background")

# --- Slide 2: What is a MoE FFN Expert? -----------------------------------
s, top = content_slide("What is a MoE FFN Expert?",
                       "Each MoE layer routes each token to K=8 of N=128 experts via a learned router")
add_text(s, CONTENT_L, top, CONTENT_W, 0.5,
         [[("Each expert is a SwiGLU FFN with three weight matrices and intermediate dimension I = 768:",
            dict(size=BODY_SZ, color=INK))]])
# equation-ish code block
add_rect(s, CONTENT_L, top + 0.62, CONTENT_W, 1.15, WHITE, line=RGBColor(0xD5,0xDE,0xE7))
add_text(s, CONTENT_L + 0.25, top + 0.74, CONTENT_W - 0.5, 1.0,
         [[("h = SiLU(W_gate · x)  ⊙  (W_up · x)", dict(size=20, name="Consolas", color=NAVY, bold=True)),
           ("      ← intermediate, dim I", dict(size=16, name="Consolas", color=MUTE))],
          [("y = W_down · h", dict(size=20, name="Consolas", color=NAVY, bold=True)),
           ("                                  ← output, dim d", dict(size=16, name="Consolas", color=MUTE))]],
         space_after=8)
bullets(s, CONTENT_L, top + 2.0, CONTENT_W, 2.2, [
    ([("gate_proj", dict(bold=True, color=TEAL)), (": produces a sparse gating signal (SiLU zeros many channels)", {})], 0),
    ([("up_proj", dict(bold=True, color=TEAL)), (": produces the value signal per channel", {})], 0),
    ([("down_proj", dict(bold=True, color=TEAL)), (": projects the intermediate back to hidden dim", {})], 0),
    ([("The element-wise product ⊙ means each channel j is an ", {}),
      ("independent computation path", dict(bold=True, color=NAVY)), (".", {})], 0),
], space_after=10)

# --- Slide 3: Why MoE Prevails --------------------------------------------
s, top = content_slide("Why MoE Prevails",
                       "10× more parameters at the same inference cost")
bullets(s, CONTENT_L, top, CONTENT_W, 2.4, [
    ([("Total params: ", {}), ("30B", dict(bold=True, color=NAVY)),
      (";  active per token: ", {}), ("~3B", dict(bold=True, color=NAVY)),
      ("  (only K=8 of 128 experts fire)", {})], 0),
    ("Quality scales with total params; latency scales with active params", 0),
    ([("MoE decouples quality from cost → ", {}),
      ("dominates dense models at equal FLOP budget", dict(bold=True, color=NAVY))], 0),
], space_after=12)
# callout box
add_rect(s, CONTENT_L, top + 2.5, CONTENT_W, 1.55, NAVY)
add_rect(s, CONTENT_L, top + 2.5, 0.14, 1.55, GOLD)
add_text(s, CONTENT_L + 0.4, top + 2.68, CONTENT_W - 0.8, 1.25,
         [[("But at decode, the bottleneck is memory bandwidth.", dict(size=22, bold=True, color=GOLD))],
          [("Even “active” params must be loaded from memory every token.  ", dict(size=BODY_SZ, color=WHITE)),
           ("Latency ∝ bytes loaded per token.", dict(size=BODY_SZ, bold=True, color=WHITE))]],
         space_after=8, line_spacing=1.05)

# --- Slide 4: What Needs to Be Improved -----------------------------------
s, top = content_slide("What Needs to Be Improved")
rows = [
    ["Goal", "Benefit", "Target"],
    ["Reduce total params", "Fit in device DRAM; less memory for serving", "Edge (L4, Cor3, phone) and cloud"],
    ["Reduce active params", "Fewer bytes loaded per token → lower latency", "Single-batch decode (bandwidth-bound)"],
]
gt, th = add_table(s, CONTENT_L, top + 0.3, CONTENT_W, rows,
                   col_w=[3.2, 5.4, 3.53], font=18, row_h=1.0, header_h=0.55,
                   col_align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
add_text(s, CONTENT_L, top + 0.3 + th + 0.45, CONTENT_W, 1.2,
         [[("Two orthogonal levers. ", dict(size=BODY_SZ, bold=True, color=NAVY)),
           ("This work targets ", dict(size=BODY_SZ, color=INK)),
           ("active", dict(size=BODY_SZ, bold=True, italic=True, color=TEAL)),
           (" params per token (the decode bottleneck), and composes with ", dict(size=BODY_SZ, color=INK)),
           ("total", dict(size=BODY_SZ, bold=True, italic=True, color=TEAL)),
           ("-param reduction via MoBE.", dict(size=BODY_SZ, color=INK))]],
         line_spacing=1.1)

# --- Section: Overview -----------------------------------------------------
section_slide("Overview of Current Design & Results")

# --- Slide 5: Framework at a Glance ---------------------------------------
s, top = content_slide("Framework at a Glance — Results First",
                       "−50% expert-FFN active parameters → −1 MMLU pt, no training")
add_text(s, CONTENT_L, top, CONTENT_W, 0.4,
         [[("Score channels per token by ", dict(size=18)),
           ("up_proj", dict(size=18, bold=True, color=TEAL)),
           (" activation magnitude; reduce ", dict(size=18)),
           ("gate_proj + down_proj", dict(size=18, bold=True, color=TEAL)),
           (".", dict(size=18))]])
rows = [
    ["Configuration", "up | gate | down", "FFN cut", "MMLU", "HellaSwag", "ARC-C", "TruthfulQA", "Avg"],
    ["Dense baseline (K=8)", "−0% | −0% | −0%", "0%", "79.6", "—", "69.7", "—", "—"],
    ["Dynamic −75%", "−0% | −75% | −75%", "−50%", "78.6", "75.4", "66.0", "51.1", "67.8"],
    ["Dynamic −87.5%", "−0% | −87.5% | −87.5%", "−58.3%", "75.3", "71.5", "63.2", "50.8", "65.2"],
]
gt, th = add_table(s, CONTENT_L, top + 0.5, CONTENT_W, rows,
                   col_w=[2.55, 2.35, 1.0, 0.98, 1.35, 0.95, 1.5, 1.45],
                   font=14, row_h=0.5, header_h=0.62, highlight_rows=(2,))
yb = top + 0.5 + th + 0.25
bullets(s, CONTENT_L, yb, CONTENT_W, 1.0, [
    ([("up_proj", dict(bold=True, color=TEAL)), (" scores channels per token via |up·x|; only top-B channels of ", {}),
      ("gate/down", dict(bold=True)), (" are activated", {})], 0),
    ([("No training, no weight changes — a pure inference-time decision", dict(bold=True, color=NAVY))], 0),
], size=16, space_after=5)
add_rect(s, CONTENT_L, CONTENT_BOT - 0.62, CONTENT_W, 0.5, RGBColor(0xEC,0xDF,0xBE))
add_text(s, CONTENT_L + 0.25, CONTENT_BOT - 0.60, CONTENT_W - 0.5, 0.46,
         [[("Edge efficiency: ", dict(size=17, color=INK)),
           ("2.14× decode speedup", dict(size=17, bold=True, color=NAVY)),
           ("  on real 30B offloaded generation (single L4 GPU).", dict(size=17, color=INK))]],
         anchor=MSO_ANCHOR.MIDDLE)

# --- Section: Motivation ---------------------------------------------------
section_slide("Motivation")

# --- Slide 6: Channel experts ---------------------------------------------
s, top = content_slide("Channel Experts", "The effective expert unit is a single intermediate channel")
add_text(s, CONTENT_L, top, CONTENT_W, 0.35,
         [[("A standard MoE FFN expert computes:", dict(size=18, color=INK))]])
add_equation(s, "eq/eq_channel_expert_1.png", top + 0.42, max_w=9.5, max_h=0.68)
add_text(s, CONTENT_L, top + 1.2, CONTENT_W, 0.35,
         [[("The j-th ", dict(size=18, color=INK)), ("channel expert", dict(size=18, bold=True, color=NAVY)),
           (" of expert e is the rank-1 computation path:", dict(size=18, color=INK))]])
add_equation(s, "eq/eq_channel_expert_2.png", top + 1.62, max_w=10.5, max_h=0.68)
bullets(s, CONTENT_L, top + 2.5, CONTENT_W, 1.4, [
    ([("Block output = ", {}), ("Σ", dict(bold=True)), (" over K·I channel experts — a sum of rank-1 paths", {})], 0),
    ([("The gating SiLU(w_gate,j·x) is a ", {}), ("soft on/off switch", dict(bold=True, color=NAVY)),
      (": near zero → the whole path contributes nothing, regardless of up/down", {})], 0),
], size=17, space_after=6)
add_rect(s, CONTENT_L, CONTENT_BOT - 0.66, CONTENT_W, 0.56, NAVY)
add_rect(s, CONTENT_L, CONTENT_BOT - 0.66, 0.14, 0.56, GOLD)
add_text(s, CONTENT_L + 0.35, CONTENT_BOT - 0.64, CONTENT_W - 0.7, 0.52,
         [[("A token doesn’t need 8 whole experts — it needs a sparse, token-specific subset of channels across them.",
            dict(size=17, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)

# --- Slide 7: A Sparse Subset Suffices ------------------------------------
s, top = content_slide("A Sparse Subset of Channels Suffices",
                       "Per token, activation magnitude concentrates on a handful of channels")
fw, fh = add_figure(s, "fig_sparse_suffices.png", top + 0.05, max_w=7.2, max_h=3.9, center_x=CONTENT_L)
tx = CONTENT_L + fw + 0.35
bullets(s, tx, top + 0.15, CONTENT_R - tx, 4.0, [
    ([("(a) Long-tailed output.", dict(bold=True, color=NAVY)), (" Histogram of SwiGLU output h_j (log y-axis).", {})], 0),
    ([("(b) Few uneven neurons.", dict(bold=True, color=NAVY)), (" After masking bottom 95%: mean 403 acts/neuron, but 8 neurons fire >5× the mean.", {})], 0),
    ([("Takeaway: ", dict(bold=True, color=TEAL)), ("top ~12.5% of channels by |h| captures a median 50% of total output magnitude (top 50% → 90%).", {})], 0),
], size=16, space_after=10)
add_text(s, CONTENT_L, top + 4.05, fw, 0.4,
         [[("Single expert (layer 0, expert 0), 8,000 WikiText-2 tokens.", dict(size=11, italic=True, color=MUTE))]],
         align=PP_ALIGN.CENTER)

# --- Slide 8: The Subset Is Token-Specific --------------------------------
s, top = content_slide("The Subset Is Token-Specific",
                       "Which channels matter is re-decided every token → no fixed subset works")
fw, fh = add_figure(s, "fig_token_specific.png", top + 0.05, max_w=7.0, max_h=3.9, center_x=CONTENT_L)
tx = CONTENT_L + fw + 0.35
bullets(s, tx, top + 0.15, CONTENT_R - tx, 3.6, [
    ([("(a)", dict(bold=True, color=NAVY)), (" Consecutive tokens routed to the same expert share only ", {}),
      ("13%", dict(bold=True, color=NAVY)), (" of their kept channels (budget ρ = 0.125).", {})], 0),
    ([("(b)", dict(bold=True, color=NAVY)), (" Almost no channel is stable: ", {}),
      ("0%", dict(bold=True, color=NAVY)), (" kept >95% of the time; only 12% kept <5%.", {})], 0),
], size=16, space_after=10)
add_rect(s, CONTENT_L, CONTENT_BOT - 0.95, CONTENT_W, 0.78, NAVY)
add_rect(s, CONTENT_L, CONTENT_BOT - 0.95, 0.14, 0.78, GOLD)
add_text(s, CONTENT_L + 0.35, CONTENT_BOT - 0.93, CONTENT_W - 0.7, 0.74,
         [[("Conclusion: the method must select online, per token, at channel granularity — offline / static baselines top out far below.",
            dict(size=16, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)

# --- Section: High-Level Framework ----------------------------------------
section_slide("High-Level Framework")

# --- Slide 9: Per-Token Channel Activation --------------------------------
s, top = content_slide("Per-Token Channel Activation",
                       "For each token, only the channels that contribute most are activated")
fw, fh = add_figure(s, "fig_channel_activation.png", top + 0.1, max_w=11.0, max_h=3.7)
bullets(s, CONTENT_L, top + fh + 0.3, CONTENT_W, 0.9, [
    ([("Fixed per-token budget: ", {}), ("B channels out of K·I total", dict(bold=True, color=NAVY))], 0),
    ([("Different tokens → different active subsets", dict(bold=True, color=NAVY))], 0),
], size=18, space_after=6)

# --- Slide 10: Design Challenges ------------------------------------------
s, top = content_slide("Design Challenges", "Two core problems")
# two numbered cards
cw = (CONTENT_W - 0.4) / 2
for i, (num, head, pts) in enumerate([
    ("1", "How to find the channels that contribute most for each token?",
     ["A learned router of size d × (N·I) would be as large as the MoE itself",
      "The existing router only knows which experts, not which channels within"]),
    ("2", "How does this bring real throughput acceleration?",
     ["If you must compute the full expert to know what to skip, there’s no saving",
      "The selection must be available before the main compute"]),
]):
    cx = CONTENT_L + i * (cw + 0.4)
    add_rect(s, cx, top + 0.1, cw, 3.7, WHITE, line=RGBColor(0xD5,0xDE,0xE7))
    add_rect(s, cx, top + 0.1, cw, 0.06, GOLD)
    add_rect(s, cx + 0.28, top + 0.4, 0.7, 0.7, NAVY)
    add_text(s, cx + 0.28, top + 0.4, 0.7, 0.7,
             [[(num, dict(size=34, bold=True, color=GOLD))]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx + 1.15, top + 0.42, cw - 1.4, 1.2,
             [[(head, dict(size=18, bold=True, color=NAVY))]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, cx + 0.3, top + 1.7, cw - 0.6, 1.9, [(p, 0) for p in pts], size=16, space_after=8)

# --- Slide 11: Our Framework ----------------------------------------------
s, top = content_slide("Our Framework", "Use up_proj output as a built-in per-channel scorer")
bullets(s, CONTENT_L, top, CONTENT_W * 0.52, 3.4, [
    ([("Compute up·x at full width", dict(bold=True, color=NAVY)), (" → score channels by activation magnitude", {})], 0),
    ([("Keep top-B channels", dict(bold=True, color=NAVY)), (" → activate only those rows of gate and columns of down", {})], 0),
    ([("up_proj storage-compressed via MoBE", dict(bold=True, color=NAVY)), (" (shared basis) to reduce its cost", {})], 0),
    ([("A predictor estimates the next layer’s mask", dict(bold=True, color=NAVY)), (" → prefetch params with no waiting", {})], 0),
], size=16, space_after=10)
add_figure(s, "fig_framework.png", top + 0.15, max_w=6.0, max_h=3.6,
           center_x=CONTENT_L + CONTENT_W * 0.52 + 0.25)

# --- Section: Details of Design --------------------------------------------
section_slide("Details of Design")

# --- Slide 12: Channel Scoring and Selection ------------------------------
s, top = content_slide("Channel Scoring and Selection",
                       "Score = per-channel contribution magnitude to the block output")
add_equation(s, "eq/eq_scoring.png", top + 0.05, max_w=8.5, max_h=0.7)
bullets(s, CONTENT_L, top + 0.95, CONTENT_W, 2.0, [
    ([("g_e", dict(bold=True, color=TEAL)), (":  router weight — puts all experts on one comparable scale", {})], 0),
    ([("|u_e,j(x)| = |w_up,j · x|", dict(bold=True, color=TEAL)), (":  activation magnitude (the token-specific signal)", {})], 0),
    ([("‖W_down[:,j]‖₂", dict(bold=True, color=TEAL)), (":  column norm (how much channel j amplifies into the output)", {})], 0),
], size=18, space_after=10)
add_rect(s, CONTENT_L, top + 2.95, CONTENT_W, 1.28, NAVY)
add_rect(s, CONTENT_L, top + 2.95, 0.14, 1.28, GOLD)
add_text(s, CONTENT_L + 0.35, top + 3.06, CONTENT_W - 0.7, 1.06,
         [[("Global top-B across all K·I channels.", dict(size=19, bold=True, color=GOLD))],
          [("Per-expert quotas emerge from a single threshold — a dominated expert may get 0 channels. No per-expert floor needed.",
            dict(size=17, color=WHITE))]], space_after=6, line_spacing=1.05)

# --- Slide 13: Why up_proj as the Channel Router --------------------------
s, top = content_slide("Why up_proj as the Channel Router",
                       "The architecture already contains a per-channel router — no extra params")
bullets(s, CONTENT_L, top, CONTENT_W, 1.3, [
    ("up·x determines the value each channel carries into the output; its magnitude predicts which channels dominate h", 0),
    ("Using up_proj to select channels lets both gate and down run at reduced width", 0),
], size=17, space_after=6)
rows = [
    ["", "Extra params", "Extra compute", "Result"],
    ["Learned router", "d × N·I  (huge)", "full matmul", "impractical"],
    ["up_proj (ours)", "0", "already computed", "gate + down reduced"],
]
gt, th = add_table(s, CONTENT_L, top + 1.35, CONTENT_W, rows,
                   col_w=[2.7, 3.0, 3.4, 3.03], font=16, row_h=0.5, header_h=0.5,
                   highlight_rows=(2,))
add_text(s, CONTENT_L, top + 1.35 + th + 0.2, CONTENT_W, 0.9,
         [[("Limitation: ", dict(size=16, bold=True, color=NAVY)),
           ("up_proj must run at full width (it IS the scoring signal); its cost is addressed by MoBE.  ", dict(size=16, color=INK))],
          [("Two scoring variants: |up·x| (shipped) and |SiLU(gate·x)| (mirror). Mirror wins PPL (11.18 vs 16.89) but −1.4 MMLU.",
            dict(size=16, color=INK))]], space_after=4, line_spacing=1.05)

# --- Slide 14: Channel Expert Predictor -----------------------------------
s, top = content_slide("Channel Expert Predictor",
                       "Predict layer i+1’s channel mask during layer i’s compute")
bullets(s, CONTENT_L, top, CONTENT_W * 0.55, 3.4, [
    ([("Problem:", dict(bold=True, color=NAVY)), (" must wait for up·x to know which channels to fetch → sequential dependency", {})], 0),
    ([("Adjacent-layer hidden states: cosine sim > 0.95", dict(bold=True)), (" (all layers except L0)", {})], 0),
    ([("Predictor: â_up^(i+1) ≈ x^(i) · W_up^(i+1)", dict(bold=True, color=TEAL)), (" — parameter-free, reuses next layer’s weights", {})], 0),
    ("Predicted mask available one layer early → prefetch gate/down from memory", 0),
], size=16, space_after=10)
rows = [
    ["", "Exact mask", "Predicted mask"],
    ["MMLU", "78.6", "77.1  (−1.5)"],
    ["PPL", "16.89", "15.11"],
    ["Latency", "serial", "parallel (overlap)"],
    ["Recall", "1.000", "0.777"],
]
tx = CONTENT_L + CONTENT_W * 0.55 + 0.3
gt, th = add_table(s, tx, top + 0.1, CONTENT_R - tx, rows,
                   col_w=[1.6, 1.55, (CONTENT_R - tx) - 3.15], font=15, row_h=0.52, header_h=0.5)
add_text(s, CONTENT_L, CONTENT_BOT - 0.5, CONTENT_W, 0.42,
         [[("Costs −1.5 MMLU pts but converts dynamic selection into a latency-free memory access pattern.",
            dict(size=17, bold=True, color=NAVY))]])

# --- Section: Experiment Results -------------------------------------------
section_slide("Experiment Results")

# --- Slide 15: Two Scoring Signals ----------------------------------------
s, top = content_slide("Two Scoring Signals — |up| vs |SiLU(gate)|",
                       "Both valid; choice depends on target metric")
rows = [
    ["Method", "up | gate | down", "FFN cut", "MMLU", "HellaSwag", "ARC-C", "TruthfulQA", "Avg"],
    ["Dense baseline", "−0|−0|−0", "0%", "79.6", "—", "69.7", "—", "—"],
    ["|up| (cut gate+down)", "−0|−75|−75", "−50%", "78.6", "75.4", "66.0", "51.1", "67.8"],
    ["|SiLU(gate)| (cut up+down)", "−75|−0|−75", "−50%", "77.2", "75.8", "67.3", "48.3", "67.2"],
    ["|up|", "−0|−87.5|−87.5", "−58.3%", "75.3", "71.5", "63.2", "50.8", "65.2"],
    ["|SiLU(gate)|", "−87.5|−0|−87.5", "−58.3%", "74.0", "72.5", "64.4", "45.1", "64.0"],
]
gt, th = add_table(s, CONTENT_L, top + 0.1, CONTENT_W, rows,
                   col_w=[3.0, 2.15, 1.05, 0.95, 1.4, 0.93, 1.55, 1.1],
                   font=14, row_h=0.44, header_h=0.55, highlight_rows=(2,))
yb = top + 0.1 + th + 0.25
bullets(s, CONTENT_L, yb, CONTENT_W, 1.1, [
    ([("|up| wins MMLU + TruthfulQA;  |SiLU(gate)| wins HellaSwag + ARC-C + PPL", dict(bold=True, color=NAVY)), (" (11.18 vs 16.89)", {})], 0),
    ([("Mechanism:", dict(bold=True, color=TEAL)), (" |up| can keep channels whose gate SiLU has closed (budget waste); |SiLU(gate)| cannot", {})], 0),
], size=16, space_after=6)

# --- Slide 16: Stacking with Top-K Reduction ------------------------------
s, top = content_slide("Stacking with Top-K Reduction",
                       "Fewer experts × narrower experts compose orthogonally")
rows = [
    ["Method", "up | gate | down", "FFN cut", "MMLU", "HellaSwag"],
    ["Dense baseline (K=8)", "−0 | −0 | −0", "0%", "79.5", "78.56"],
    ["Top-4 only", "−0 | −0 | −0  (×4/8)", "−50%", "77.4", "75.96"],
    ["Dynamic −75% only (K=8)", "−0 | −75 | −75", "−50%", "78.6", "75.4"],
    ["Top-4 + dynamic −50%", "−0 | −50 | −50  (×4/8)", "−66.7%", "77.2", "74.0"],
    ["Top-4 + dynamic −75%", "−0 | −75 | −75  (×4/8)", "−75.0%", "74.3", "70.0"],
]
gt, th = add_table(s, CONTENT_L, top + 0.1, CONTENT_W, rows,
                   col_w=[3.5, 3.4, 1.6, 1.6, 2.03],
                   font=15, row_h=0.46, header_h=0.5, highlight_rows=(4, 5))
yb = top + 0.1 + th + 0.25
bullets(s, CONTENT_L, yb, CONTENT_W, 1.1, [
    ("At −50% FFN, narrowing alone (78.6) outperforms expert-dropping (77.4)", 0),
    ("At −75% FFN, stacking retains knowledge expert-dropping destroys (reduce-top-k 8→2 collapses to HellaSwag 49.4)", 0),
    ([("Orthogonal:", dict(bold=True, color=TEAL)), (" top-k cuts whole experts; dynamic cuts channels within surviving experts", {})], 0),
], size=15, space_after=5)

# --- Slide 17: Predictor Accuracy Cost ------------------------------------
s, top = content_slide("Channel Expert Predictor — Accuracy Cost",
                       "Predicting the mask one layer early costs −1.5 MMLU pts (−50% FFN, |up|)")
rows = [
    ["Configuration", "up | gate | down", "MMLU", "HellaSwag", "ARC-C", "TruthfulQA", "Avg"],
    ["Exact mask", "−0 | −75 | −75", "78.6", "75.4", "66.0", "51.1", "67.8"],
    ["Predicted mask (FloE)", "−0 | −75 | −75", "77.1", "75.2", "65.0", "52.1", "67.4"],
]
gt, th = add_table(s, CONTENT_L, top + 0.2, CONTENT_W, rows,
                   col_w=[3.0, 2.5, 1.35, 1.6, 1.28, 1.7, 0.7],
                   font=15, row_h=0.55, header_h=0.6, highlight_rows=(2,))
bullets(s, CONTENT_L, top + 0.2 + th + 0.35, CONTENT_W, 1.6, [
    ([("Measured recall: ", {}), ("0.777", dict(bold=True, color=NAVY)), ("  (vs FloE’s ~0.95 on Mixtral-8×7B)", {})], 0),
    ([("Average cost: ", {}), ("−0.4 pts", dict(bold=True, color=NAVY)), (" across 4 tasks; only MMLU loses measurably (−1.5)", {})], 0),
    ("Parameter-free; enables the latency-hiding prefetch", 0),
], size=18, space_after=10)

# --- Slide 18: Efficiency — Edge Offload ----------------------------------
s, top = content_slide("Efficiency — Edge Offload (Real 30B, Single L4)",
                       "Experts in CPU DRAM, AIME-24 generation, batch-1 decode")
rows = [
    ["Variant", "ms/tok", "tok/s", "vs dense", "Peak GPU"],
    ["Dense (all experts offloaded)", "526", "1.90", "1.00×", "3.2 GB"],
    ["Dynamic (−75% gate+down)", "291", "3.43", "1.81×", "3.2 GB"],
    ["+ Predicted prefetch", "246", "4.01", "2.14×", "3.2 GB"],
]
gt, th = add_table(s, CONTENT_L, top + 0.1, CONTENT_W, rows,
                   col_w=[4.3, 1.7, 1.7, 1.85, 2.58],
                   font=16, row_h=0.5, header_h=0.5, highlight_rows=(2, 3))
bullets(s, CONTENT_L, top + 0.1 + th + 0.3, CONTENT_W, 1.6, [
    ([("Halving PCIe bytes buys ", {}), ("1.81×", dict(bold=True, color=NAVY)), (" on the clock (37.75 vs 75.50 MB/layer)", {})], 0),
    ([("Predicted prefetch overlaps next-layer fetch with compute → ", {}), ("2.14× total", dict(bold=True, color=NAVY))], 0),
    ("Peak GPU only 3.2 GB — 30B on one L4 with 20 GB to spare; prefill ~28 tok/s (mask-union saturates)", 0),
], size=16, space_after=7)

# --- Slide 19: Efficiency — Cloud Resident --------------------------------
s, top = content_slide("Efficiency — Cloud Resident (4× L4)",
                       "Experts resident in HBM, AIME-24 generation")
rows = [
    ["Variant", "Batch", "Prefill (tok/s)", "Decode (ms/tok)", "Peak GPU"],
    ["Dense", "1", "172", "120", "15.7 GB"],
    ["Dense", "4", "584", "157", "15.9 GB"],
    ["Dynamic (gathered)", "1", "125", "410", "17.4 GB"],
    ["Dynamic (gathered)", "4", "412", "434", "17.6 GB"],
]
gt, th = add_table(s, CONTENT_L, top, CONTENT_W, rows,
                   col_w=[3.6, 1.6, 2.8, 2.8, 1.333],
                   font=15, row_h=0.4, header_h=0.46)
yb = top + th + 0.18
add_text(s, CONTENT_L, yb, CONTENT_W, 0.55,
         [[("Cloud resident with MoBE-factored up is a regression (2.38× read amplification); without MoBE the gathered kernel gives ", dict(size=15, color=INK)),
           ("1.34× per-layer", dict(size=15, bold=True, color=NAVY)),
           (".  Multi-batch decode dilutes per-token sparsity (mask-union saturates).", dict(size=15, color=INK))]],
         line_spacing=1.02)
# summary strip
add_text(s, CONTENT_L, yb + 0.72, CONTENT_W, 0.35,
         [[("Summary — two deployment targets", dict(size=17, bold=True, color=TEAL))]])
rows2 = [
    ["Setting", "Bottleneck", "Decode benefit", "Takeaway"],
    ["Edge (offload, batch-1)", "PCIe", "2.14× (4.0 tok/s)", "Primary target; prefetch essential"],
    ["Cloud (resident, batch≥1)", "HBM", "1.34× per-layer", "Modest; mask-union saturates at batch>1"],
]
gt2, th2 = add_table(s, CONTENT_L, yb + 1.12, CONTENT_W, rows2,
                     col_w=[3.3, 1.7, 2.7, 4.433], font=14, row_h=0.46, header_h=0.46,
                     highlight_rows=(1,), col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT])

# --- Section: Failing Experiments ------------------------------------------
section_slide("Failing Experiments — What Didn’t Work")

# --- Slide 20: Offline Static Channel Ranking -----------------------------
s, top = content_slide("Offline Static Channel Ranking (Level 1)",
                       "A fixed per-expert channel order, budget-agnostic")
bullets(s, CONTENT_L, top, CONTENT_W, 1.5, [
    ([("Idea:", dict(bold=True, color=NAVY)), (" precompute a fixed per-expert channel order; keep top-B at inference. No full-width compute at runtime — gate and up both run reduced.", {})], 0),
    ([("Method:", dict(bold=True, color=NAVY)), (" Pivoted Cholesky on activation×weight coupling → nested priority order. Online: score by g²·σ, keep global top-B.", {})], 0),
], size=17, space_after=8)
rows = [
    ["Reduction", "Offline Level 1", "Reduce top-k (8→k)"],
    ["−37.5%", "76.30", "77.1  (8→5)"],
    ["−50%", "74.26", "75.2  (8→4)"],
    ["−62.5%", "70.54", "69.8  (8→3)"],
    ["−75%", "63.60", "49.4  (8→2)"],
]
gt, th = add_table(s, CONTENT_L, top + 1.7, CONTENT_W * 0.55, rows,
                   col_w=[2.2, 2.4, (CONTENT_W * 0.55) - 4.6], font=16, row_h=0.44, header_h=0.5)
tx = CONTENT_L + CONTENT_W * 0.55 + 0.35
add_text(s, tx, top + 1.9, CONTENT_R - tx, 2.5,
         [[("Level 1 dominates expert-dropping at deep cuts (−62.5%, −75%)…", dict(size=18, color=INK))],
          [("…but is still ", dict(size=18, color=INK)),
           ("4–15 pts below online selection.", dict(size=18, bold=True, color=NAVY))]],
         space_after=10, line_spacing=1.1)

# --- Slide 21: Why Offline Methods Fail -----------------------------------
s, top = content_slide("Why Offline Methods Fail",
                       "Channel contribution differs across tokens → fixed ranking is fundamentally limited")
tbl_w = 5.1
rows = [
    ["Reduction", "Best offline", "Online", "Gap"],
    ["−50%", "74.26", "78.54", "4.3"],
    ["−75%", "63.60", "78.28", "14.7"],
    ["−87.5%", "44.15", "76.84", "32.7"],
]
gt, th = add_table(s, CONTENT_L, top, tbl_w, rows,
                   col_w=[1.5, 1.35, 1.15, 1.1], font=14, row_h=0.46,
                   header_h=0.5, highlight_rows=(3,))
fx = CONTENT_L + tbl_w + 0.35
add_figure(s, "fig_fixed_fails.png", top - 0.05, max_w=CONTENT_R - fx, max_h=2.55, center_x=fx)
bullets(s, CONTENT_L, top + max(th, 2.55) + 0.2, CONTENT_W, 1.7, [
    ([("(a)", dict(bold=True, color=NAVY)), (" Consecutive same-expert tokens share only 7–20% of kept channels (layer 24, ρ=0.125).", {})], 0),
    ([("(b)", dict(bold=True, color=NAVY)), (" One token needs 12.5% of channels — but the union over a prefill reaches 35% by 64 tokens, 74% by 2048.", {})], 0),
    ([("99.3% of channels have variable utility; the entire headroom above offline is ", {}),
      ("per-token activation information.", dict(bold=True, color=TEAL))], 0),
], size=15, space_after=5)

# --- Section: Future Steps -------------------------------------------------
section_slide("Future Steps")

# --- Slide 22: MoBE — Reduce Total Parameters -----------------------------
s, top = content_slide("MoBE — Reduce Total Parameters",
                       "Factorize expert weights into shared basis + per-expert transform")
fw, fh = add_figure(s, "mobe_mechanism.png", top - 0.05, max_w=CONTENT_W, max_h=1.95)
eq_y = top + fh + 0.05
add_equation(s, "eq/eq_mobe.png", eq_y, max_w=4.9, max_h=0.72, center=False, x=CONTENT_L)
bullets(s, CONTENT_L + 5.35, eq_y - 0.05, CONTENT_R - (CONTENT_L + 5.35), 1.2, [
    ([("Bⱼ ∈ ℝ ʳˣᵈ", dict(bold=True, color=TEAL)), (": m shared basis matrices (stored once for all 128 experts)", {})], 0),
    ([("Aₑ ∈ ℝ ᵖˣʳ", dict(bold=True, color=TEAL)), (": per-expert transform.  f = SiLU (weight-space nonlinearity)", {})], 0),
    ([("Data-free fit: Adam, lr = 0.07, 2000 steps per (layer, type)", dict(color=INK))], 0),
], size=13, space_after=4)
rows = [
    ["Setting", "MoE storage ↓", "MMLU", "HellaSwag", "PPL"],
    ["Baseline", "0%", "79.6", "77.68", "8.70"],
    ["MoBE up50% (m=16, r=768)", "−16.7%", "76.6", "—", "12.91"],
    ["MoBE even-split −33% (m=38)", "−33%", "76.83", "73.13", "10.10"],
]
gt, th = add_table(s, CONTENT_L, eq_y + 1.0, CONTENT_W, rows,
                   col_w=[4.2, 2.4, 1.7, 2.0, 1.833], font=14, row_h=0.38, header_h=0.42,
                   highlight_rows=(3,))

# --- Slide 23: MoBE + Dynamic ---------------------------------------------
s, top = content_slide("MoBE + Dynamic — Stacking Two Orthogonal Axes",
                       "MoBE reduces storage; dynamic reduces active. The two compose.")
rows = [
    ["Configuration", "Storage ↓", "Active FFN cut", "MMLU", "HellaSwag", "PPL"],
    ["Dense baseline", "0%", "0%", "79.6", "78.56", "10.89"],
    ["MoBE up50 alone", "−16.7%", "−16.7%", "76.6", "—", "12.91"],
    ["Dynamic −75% alone (K=8)", "0%", "−50%", "78.6", "75.4", "16.89"],
    ["MoBE up50 + dynamic −75%", "−16.7%", "−66.7%", "76.3", "71.3", "18.82"],
]
gt, th = add_table(s, CONTENT_L, top + 0.15, CONTENT_W, rows,
                   col_w=[3.7, 1.7, 2.3, 1.3, 1.83, 1.3], font=15, row_h=0.46, header_h=0.5,
                   highlight_rows=(4,))
bullets(s, CONTENT_L, top + 0.15 + th + 0.3, CONTENT_W, 1.5, [
    ("MoBE compresses up_proj storage (the scoring signal); dynamic cuts gate/down active per token", 0),
    ([("Composed −66.7% active cut (MMLU 76.3) dominates pushing either axis alone to the same depth", dict(bold=True, color=NAVY))], 0),
    ("On edge, MoBE also reduces DRAM footprint (48.3 GB vs 58.0 GB), fitting smaller devices", 0),
], size=16, space_after=7)

# --- Slide 24: Next Steps --------------------------------------------------
s, top = content_slide("Next Steps")
cw = (CONTENT_W - 0.6) / 3
cols = [
    ("System design for throughput", [
        "Fused Triton kernel: up[full] → select M → gate[M]·x → SiLU → ⊙ up[M] → down[M], one launch",
        "End-to-end offload runtime with double-buffered prefetch → interactive 30B on edge",
        "Memory-fetch pattern optimization for the multi-layer pipeline",
    ]),
    ("Reduce total parameters further", [
        "MoBE beats Nyström at 1.5× compression (73.13 HS vs 65.10 at −33%)",
        "Stack on Nyström-compressed base: another 1.5× on top (running now)",
    ]),
    ("Learn per-token expert budget", [
        "Current: fixed K=8 experts/token. Many tokens need fewer",
        "Combine with top-p routing → dynamic total budget per token (not just channel allocation)",
    ]),
]
for i, (head, pts) in enumerate(cols):
    cx = CONTENT_L + i * (cw + 0.3)
    add_rect(s, cx, top, cw, 4.1, WHITE, line=RGBColor(0xD5,0xDE,0xE7))
    add_rect(s, cx, top, cw, 0.72, NAVY)
    add_text(s, cx + 0.18, top, cw - 0.36, 0.72,
             [[(head, dict(size=16, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, cx + 0.2, top + 0.9, cw - 0.4, 3.0, [(p, 0) for p in pts], size=14, space_after=8)

# --- Slide 25: Summary -----------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT)
add_rect(s, 0, 0, 0.35, SH, GOLD)
add_text(s, MARGIN + 0.2, 0.45, SW - 2 * MARGIN, 0.9,
         [[("Summary", dict(size=36, bold=True, color=NAVY))]])
add_text(s, MARGIN + 0.2, 1.25, SW - 2 * MARGIN, 0.5,
         [[("Per-token dynamic channel selection: the right granularity for MoE efficiency",
            dict(size=20, bold=True, color=TEAL))]])
pts = [
    ([("Channel experts", dict(bold=True, color=TEAL)), ("  —  the effective expert unit is a single intermediate channel, not a whole expert", dict(color=INK))]),
    ([("up_proj as built-in router", dict(bold=True, color=TEAL)), ("  —  no extra params; activation magnitude scores channels per token", dict(color=INK))]),
    ([("−50% expert-FFN active → MMLU 78.6 (−1 pt), no training", dict(bold=True, color=TEAL)), ("  —  exact at budget, composes with MoBE", dict(color=INK))]),
    ([("2.14× edge decode speedup", dict(bold=True, color=TEAL)), ("  —  predicted prefetch at 97% PCIe peak; payoff on bandwidth-starved hardware", dict(color=INK))]),
]
tb = s.shapes.add_textbox(Emu(int((MARGIN + 0.2) * 914400)), Emu(int(2.05 * 914400)),
                          Emu(int((SW - 2 * MARGIN - 0.2) * 914400)), Emu(int(3.4 * 914400)))
tf = tb.text_frame; tf.word_wrap = True
for i, runs in enumerate(pts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(14); p.line_spacing = 1.05
    pPr = p._p.get_or_add_pPr()
    # numbered
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": FONT})
    buAutoNum = pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"})
    buClr = pPr.makeelement(qn("a:buClr"), {}); srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "1C7293"}); buClr.append(srgb)
    pPr.append(buClr); pPr.append(buFont); pPr.append(buAutoNum)
    pPr.set("indent", str(-int(0.35 * 914400))); pPr.set("marL", str(int(0.35 * 914400)))
    for text, kw in ([runs] if isinstance(runs, tuple) else runs):
        r = p.add_run(); r.text = text
        merged = dict(size=19); merged.update(kw); _set_font(r, **merged)
add_rect(s, MARGIN + 0.2, 5.72, SW - 2 * MARGIN - 0.2, 0.82, NAVY)
add_rect(s, MARGIN + 0.2, 5.72, 0.14, 0.82, GOLD)
add_text(s, MARGIN + 0.55, 5.80, SW - 2 * MARGIN - 0.9, 0.66,
         [[("Training-free, exact at budget, and it addresses the real bottleneck: memory bandwidth at decode.",
            dict(size=18, bold=True, color=WHITE))]], anchor=MSO_ANCHOR.MIDDLE)
add_logo(s, x=0.55)   # clear the gold spine

prs.save(str(OUT))
print(f"Saved {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
