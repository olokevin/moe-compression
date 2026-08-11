#!/usr/bin/env python3
"""Build the broader-audience intro slides as a standalone deck.

Two new slides, same design system as ``build_pptx.py`` (navy/gold palette,
Times New Roman), meant to be shown *before* the midpoint deck:

  1. Dense Model vs Mixture-of-Experts — side-by-side block diagrams
     (sourced from the GShard figure on the Hugging Face MoE blog).
  2. The Challenge: MoE on Edge — a 3-tier memory hierarchy (DRAM / SRAM /
     compute) with a per-layer fetch-vs-compute latency comparison showing
     decode is memory-bound, then the two optimization challenges.

The original ``Yequan_26_Midpoint.pptx`` is NOT modified.
Output: ``Yequan_26_Intro.pptx``.

Appended after the motivation slides is the Level-1 (offline static channel
ranking) block — the "failing experiments" story plus bonus results — sourced
from ``midpoint_level1.md`` and rendered in the same navy/gold design system.

Run from ``docs/presentation``:  python build_intro_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageChops, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
FIGS = HERE / "figs"
OUT = HERE / "Yequan_26_Intro.pptx"

# ---------------------------------------------------------------------------
# Palette (identical to build_pptx.py)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x2A, 0x47)
NAVY2 = RGBColor(0x14, 0x3A, 0x5E)
GOLD = RGBColor(0xE0, 0xA3, 0x2E)
TEAL = RGBColor(0x1C, 0x72, 0x93)
INK = RGBColor(0x22, 0x2A, 0x33)
MUTE = RGBColor(0x5C, 0x6B, 0x7A)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROWALT = RGBColor(0xEA, 0xEF, 0xF4)
RED = RGBColor(0xB4, 0x32, 0x32)
HEADFILL = NAVY
HILITE = RGBColor(0xFC, 0xF3, 0xDC)    # soft gold row/col highlight

FONT = "Times New Roman"

# tuples for PIL (RGB)
P_NAVY = (0x0E, 0x2A, 0x47)
P_NAVY2 = (0x14, 0x3A, 0x5E)
P_GOLD = (0xE0, 0xA3, 0x2E)
P_TEAL = (0x1C, 0x72, 0x93)
P_INK = (0x22, 0x2A, 0x33)
P_MUTE = (0x5C, 0x6B, 0x7A)
P_LIGHT = (0xF4, 0xF6, 0xF9)
P_WHITE = (0xFF, 0xFF, 0xFF)
P_RED = (0xB4, 0x32, 0x32)
P_GREEN = (0x2E, 0x8B, 0x3D)
P_LINE = (0xC7, 0xD1, 0xDB)

SW, SH = 13.333, 7.5
MARGIN = 0.6
CONTENT_TOP = 1.40
CONTENT_BOT = 6.82
CONTENT_L = MARGIN
CONTENT_R = SW - MARGIN
CONTENT_W = CONTENT_R - CONTENT_L

LOGO = FIGS / "amazon_logo.png"
LOGO_X, LOGO_Y, LOGO_W, LOGO_H = 0.21, 6.90, 1.60, 0.482

BODY_SZ = 20
TITLE_SZ = 30
SUBTITLE_SZ = 20

prs = Presentation()
prs.slide_width = Emu(int(SW * 914400))
prs.slide_height = Emu(int(SH * 914400))
BLANK = prs.slide_layouts[6]

_slide_no = 0


# ---------------------------------------------------------------------------
# Low-level helpers (identical to build_pptx.py)
# ---------------------------------------------------------------------------
def _set_font(run, size=BODY_SZ, bold=False, italic=False, color=INK, name=FONT):
    f = run.font
    f.name = name
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x * 914400)),
                                Emu(int(y * 914400)), Emu(int(w * 914400)),
                                Emu(int(h * 914400)))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space_after=6, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                  Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(int(0.05 * 914400))
    tf.margin_right = Emu(int(0.05 * 914400))
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for text, kw in para:
            r = p.add_run()
            r.text = text
            _set_font(r, **kw)
    return tb


def add_logo(slide, x=LOGO_X):
    slide.shapes.add_picture(str(LOGO), Emu(int(x * 914400)),
                             Emu(int(LOGO_Y * 914400)), Emu(int(LOGO_W * 914400)),
                             Emu(int(LOGO_H * 914400)))


def add_footer(slide, title_short):
    global _slide_no
    _slide_no += 1
    add_logo(slide)
    add_rect(slide, SW - 1.15, SH - 0.5, 0.12, 0.26, GOLD)
    add_text(slide, SW - 0.95, SH - 0.54, 0.75, 0.35,
             [[(str(_slide_no), dict(size=12, color=MUTE, bold=True))]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 2.0, SH - 0.52, 9.0, 0.35,
             [[("Per-Token Adaptive Channel Activation for MoE", dict(size=11, color=MUTE, italic=True))]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def content_slide(title, subtitle=None):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, LIGHT)
    add_rect(s, MARGIN, 0.34, 0.14, 0.62, GOLD)
    add_text(s, MARGIN + 0.30, 0.22, SW - 2 * MARGIN - 0.30, 0.86,
             [[(title, dict(size=TITLE_SZ, bold=True, color=NAVY))]],
             anchor=MSO_ANCHOR.MIDDLE)
    top = CONTENT_TOP
    if subtitle:
        add_text(s, MARGIN, 1.18, CONTENT_W, 0.5,
                 [[(subtitle, dict(size=SUBTITLE_SZ, bold=True, color=TEAL))]],
                 anchor=MSO_ANCHOR.TOP)
        top = 1.82
    add_footer(s, title)
    return s, top


def bullets(slide, x, y, w, h, items, size=BODY_SZ, space_after=8, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Emu(int(x * 914400)), Emu(int(y * 914400)),
                                  Emu(int(w * 914400)), Emu(int(h * 914400)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if len(item) == 3:
            content, level, kw = item
        else:
            content, level = item
            kw = {}
        p.level = level
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "▪" if level == 0 else "–"})
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": FONT})
        buClr = pPr.makeelement(qn("a:buClr"), {})
        srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "E0A32E" if level == 0 else "1C7293"})
        buClr.append(srgb)
        pPr.append(buClr)
        pPr.append(buFont)
        pPr.append(buChar)
        indent = 0.28
        pPr.set("indent", str(-int(indent * 914400)))
        pPr.set("marL", str(int((indent + level * 0.3) * 914400)))
        runs = content if isinstance(content, list) else [(content, {})]
        for text, rkw in runs:
            r = p.add_run()
            r.text = text
            merged = dict(size=size, color=INK)
            merged.update(kw)
            merged.update(rkw)
            _set_font(r, **merged)
    return tb


def add_picture_hw(slide, png_path, x, y, w, h):
    slide.shapes.add_picture(str(png_path), Emu(int(x * 914400)),
                             Emu(int(y * 914400)), Emu(int(w * 914400)),
                             Emu(int(h * 914400)))


def add_figure(slide, png_path, y, max_w=CONTENT_W, max_h=3.5, center_x=None):
    im = Image.open(png_path)
    ar = im.width / im.height
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    x = (SW - w) / 2 if center_x is None else center_x
    slide.shapes.add_picture(str(png_path), Emu(int(x * 914400)),
                             Emu(int(y * 914400)), Emu(int(w * 914400)),
                             Emu(int(h * 914400)))
    return w, h


def section_slide(title, subtitle=None):
    """Full-page section divider (navy field, gold marker) — matches build_pptx."""
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_rect(s, MARGIN, 2.75, 1.4, 0.14, GOLD)
    add_text(s, MARGIN, 3.12, SW - 2 * MARGIN, 1.2,
             [[(title, dict(size=40, bold=True, color=WHITE))]],
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    if subtitle:
        add_text(s, MARGIN, 4.05, SW - 2 * MARGIN, 0.9,
                 [[(subtitle, dict(size=20, color=RGBColor(0xC9, 0xD6, 0xE3)))]],
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_logo(s)
    return s


def add_table(slide, x, y, w, rows, col_w=None, header=True, font=14,
              highlight_rows=(), highlight_cols=(), col_align=None,
              row_h=0.34, header_h=0.46, bold_cells=()):
    """rows: list of list of cell-strings (row 0 = header if header=True).

    ``highlight_cols`` shades whole columns gold (the "winning" method columns);
    ``bold_cells`` is a set of (i, j) to embolden individually.
    """
    from pptx.enum.shapes import MSO_SHAPE  # noqa: F401  (kept for parity)
    nrow = len(rows)
    ncol = len(rows[0])
    total_h = header_h + (nrow - 1) * row_h if header else nrow * row_h
    gt = slide.shapes.add_table(nrow, ncol, Emu(int(x * 914400)), Emu(int(y * 914400)),
                                Emu(int(w * 914400)), Emu(int(total_h * 914400)))
    tbl = gt.table
    tbl.first_row = header
    tbl.horz_banding = False
    if col_w:
        for j, cw in enumerate(col_w):
            tbl.columns[j].width = Emu(int(cw * 914400))
    for i in range(nrow):
        tbl.rows[i].height = Emu(int((header_h if (header and i == 0) else row_h) * 914400))
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = Emu(int(0.06 * 914400))
            cell.margin_right = Emu(int(0.06 * 914400))
            cell.margin_top = Emu(int(0.01 * 914400))
            cell.margin_bottom = Emu(int(0.01 * 914400))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if header and i == 0:
                fill_c = HEADFILL
            elif i in highlight_rows:
                fill_c = HILITE
            elif j in highlight_cols:
                fill_c = HILITE
            elif i % 2 == 0:
                fill_c = ROWALT
            else:
                fill_c = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_c
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            if col_align:
                p.alignment = col_align[j]
            else:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(val)
            is_head = header and i == 0
            emphasize = ((i, j) in bold_cells or (i in highlight_rows and j == 0)
                         or (j in highlight_cols and i != 0))
            _set_font(r, size=font, bold=is_head or emphasize,
                      color=WHITE if is_head else INK)
    return gt, total_h


# ---------------------------------------------------------------------------
# Fonts for PIL figure generation
# ---------------------------------------------------------------------------
def _font(size, bold=False):
    names = (["DejaVuSans-Bold.ttf"] if bold else ["DejaVuSans.ttf"])
    roots = [
        "/usr/share/fonts/truetype/dejavu/",
        "/usr/share/fonts/dejavu-sans-fonts/",
        "/usr/share/fonts/dejavu/",
        "/usr/share/fonts/gnu-free/",
    ]
    for r in roots:
        for n in names:
            p = Path(r) / n
            if p.exists():
                return ImageFont.truetype(str(p), size)
    return ImageFont.load_default()


def _autocrop(im, pad=8, bg=(255, 255, 255)):
    rgb = im.convert("RGB")
    diff = ImageChops.difference(rgb, Image.new("RGB", rgb.size, bg))
    bbox = diff.getbbox()
    if bbox:
        l, t, r, b = bbox
        return im.crop((max(0, l - pad), max(0, t - pad),
                        min(im.size[0], r + pad), min(im.size[1], b + pad)))
    return im


# ---------------------------------------------------------------------------
# Prepare the dense / MoE block diagrams from the downloaded GShard figure
# ---------------------------------------------------------------------------
def prepare_block_figures():
    src = FIGS / "web" / "02_moe_block.png"
    dense_out = FIGS / "fig_intro_dense_block.png"
    moe_out = FIGS / "fig_intro_moe_block.png"
    if not src.exists():
        print(f"  WARNING: {src} missing — run the download step first")
        return dense_out, moe_out
    im = Image.open(src).convert("RGB")
    dense = _autocrop(im.crop((30, 150, 228, 525)))
    moe = _autocrop(im.crop((300, 60, 488, 648)))
    dense.save(str(dense_out))
    moe.save(str(moe_out))
    print(f"  block figures: dense {dense.size}, moe {moe.size}")
    return dense_out, moe_out


# ---------------------------------------------------------------------------
# Memory-hierarchy + latency figure (drawn with PIL)
# ---------------------------------------------------------------------------
def draw_memory_hierarchy(path: Path):
    S = 2  # supersample for crisp text
    W, H = 1020 * S, 672 * S
    img = Image.new("RGB", (W, H), P_WHITE)
    d = ImageDraw.Draw(img)

    def F(sz, bold=False):
        return _font(sz * S, bold)

    def rr(x0, y0, x1, y1, r, **kw):
        d.rounded_rectangle([x0 * S, y0 * S, x1 * S, y1 * S], radius=r * S, **kw)

    def tx(x, y, s, fill, fnt, anchor="lt"):
        d.text((x * S, y * S), s, fill=fill, font=fnt, anchor=anchor)

    def downarrow(x, y_top, y_bot, fill, width=6):
        d.line([(x * S, y_top * S), (x * S, y_bot * S)], fill=fill, width=width * S)
        d.polygon([((x - 11) * S, (y_bot - 14) * S), ((x + 11) * S, (y_bot - 14) * S),
                   (x * S, y_bot * S)], fill=fill)

    # geometry: three EQUAL-SIZED rectangles, centred, stacked top->bottom
    RX0, RX1 = 320, 700           # rectangle x-span (width 380)
    cx = (RX0 + RX1) / 2
    RH = 100                      # rectangle height (all identical)
    GAP = 72                      # arrow gap between rectangles (room for a label)
    Y0 = 76                       # top of first rectangle

    tx(W / (2 * S), 12,
       "One MoE layer at decode  (batch = 1, one token, 75.5 MB active)",
       P_NAVY, F(23, True), "mt")

    # tiers top->bottom = DRAM -> VRAM -> tensor cores
    tiers = [
        (P_NAVY, "Host DRAM", "total 30B = 60 GB"),
        (P_TEAL, "GPU HBM (VRAM)", "active K = 8 = 75.5 MB"),
        (P_GOLD, "Tensor cores", "compute — 75.5 MFLOP"),
    ]
    ys = []
    for i, (fill, title, sub) in enumerate(tiers):
        y0 = Y0 + i * (RH + GAP)
        y1 = y0 + RH
        ys.append((y0, y1))
        rr(RX0, y0, RX1, y1, 12, fill=fill)
        tw = "black" if fill == P_GOLD else P_WHITE
        subc = P_INK if fill == P_GOLD else (225, 233, 240)
        ymid = (y0 + y1) / 2
        tx(cx, ymid - 17, title, tw, F(28, True), "mm")
        tx(cx, ymid + 21, sub, subc, F(17), "mm")

    # transfer-cost labels sit in the arrow gaps between the boxes
    edges = [
        (ys[0][1], ys[1][0], "DRAM → VRAM:  5.62 ms", P_RED),
        (ys[1][1], ys[2][0], "VRAM → tensor cores:  1.19 ms", P_TEAL),
    ]
    for y_top, y_bot, label, col in edges:
        downarrow(cx, y_top + 4, y_bot - 4, P_MUTE, 6)
        gy = (y_top + y_bot) / 2
        tx(cx + 28, gy, label, col, F(22, True), "lm")

    # conclusion band
    by = ys[2][1] + 40
    rr(40, by, W / S - 40, by + 58, 10, fill=P_NAVY)
    rr(40, by, 54, by + 58, 4, fill=P_GOLD)
    tx(72, by + 29,
       "Decode is memory-bound: fetching weights (5.62 ms) dwarfs using them (1.19 ms).",
       P_GOLD, F(18, True), "lm")

    # footnote
    fy = by + 74
    tx(40, fy,
       "Per token, all 48 layers:  PCIe-offload ≈ 270 ms   vs   HBM-resident ≈ 57 ms.",
       P_INK, F(15, True), "lt")
    tx(40, fy + 24,
       "Measured on a single NVIDIA L4 (bf16); the DRAM ≪ VRAM ordering is what "
       "transfers to cloud GPUs.",
       P_MUTE, F(13), "lt")

    img = img.resize((W // S, H // S), Image.LANCZOS)
    img = _autocrop(img, pad=6)
    img.save(str(path))
    print(f"  memory hierarchy figure: {img.size}")


# ---------------------------------------------------------------------------
# Build figures
# ---------------------------------------------------------------------------
print("Preparing figures...")
DENSE_FIG, MOE_FIG = prepare_block_figures()
HIER_FIG = FIGS / "fig_intro_memory_hierarchy.png"
draw_memory_hierarchy(HIER_FIG)


# ===========================================================================
# SLIDE 1 — Dense vs MoE
# ===========================================================================
s, top = content_slide("Dense Model  vs  Mixture-of-Experts (MoE)",
                       "Two ways to build a transformer — MoE is why modern LLMs scale")

# Two block diagrams at a shared scale so the boxes are the same size.
dense_im = Image.open(DENSE_FIG)
moe_im = Image.open(MOE_FIG)
FIG_H = 3.25                                   # MoE (taller) height in inches
scale = FIG_H / moe_im.height                  # inches per source pixel (shared)
moe_w, moe_h = moe_im.width * scale, moe_im.height * scale
dense_w, dense_h = dense_im.width * scale, dense_im.height * scale

diag_top = top + 0.40
# Both diagrams sit near the centre with a "vs" chip between them; the
# descriptive text is pushed to the outer edges so nothing overlaps the chip.
lc = CONTENT_L + CONTENT_W * 0.34         # dense diagram centre
rc = CONTENT_L + CONTENT_W * 0.66         # MoE diagram centre
dense_x = lc - dense_w / 2
moe_x = rc - moe_w / 2
add_picture_hw(s, DENSE_FIG, dense_x, diag_top, dense_w, dense_h)
add_picture_hw(s, MOE_FIG, moe_x, diag_top, moe_w, moe_h)

# titles above each diagram (centred on the diagram, not the half)
add_text(s, lc - 2.2, top, 4.4, 0.4,
         [[("Dense Transformer", dict(size=20, bold=True, color=NAVY))]],
         align=PP_ALIGN.CENTER)
add_text(s, rc - 2.2, top, 4.4, 0.4,
         [[("Mixture-of-Experts", dict(size=20, bold=True, color=NAVY))]],
         align=PP_ALIGN.CENTER)

# "vs" chip centred in the gap between the two diagrams
vs_cx = (dense_x + dense_w + moe_x) / 2
add_rect(s, vs_cx - 0.32, diag_top + 1.3, 0.64, 0.5, GOLD)
add_text(s, vs_cx - 0.32, diag_top + 1.3, 0.64, 0.5,
         [[("vs", dict(size=22, bold=True, color=NAVY))]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# dense annotation — far left, right-aligned so it hugs the dense diagram
add_text(s, CONTENT_L, diag_top + 0.55, dense_x - CONTENT_L - 0.2, 2.6,
         [[("Every token flows", dict(size=15, color=INK))],
          [("through the same", dict(size=15, color=INK))],
          [("single FFN.", dict(size=15, color=INK))],
          [("", dict(size=6))],
          [("100% of params", dict(size=15, bold=True, color=TEAL))],
          [("active per token.", dict(size=15, bold=True, color=TEAL))]],
         align=PP_ALIGN.RIGHT, space_after=2, line_spacing=1.05)

# MoE annotation — far right, left-aligned so it hugs the MoE diagram
add_text(s, moe_x + moe_w + 0.2, diag_top + 0.55,
         CONTENT_R - (moe_x + moe_w) - 0.2, 3.0,
         [[("A router sends each", dict(size=15, color=INK))],
          [("token to K of N", dict(size=15, color=INK))],
          [("expert FFNs.", dict(size=15, color=INK))],
          [("", dict(size=6))],
          [("Many experts,", dict(size=15, bold=True, color=TEAL))],
          [("few fire per token.", dict(size=15, bold=True, color=TEAL))]],
         align=PP_ALIGN.LEFT, space_after=2, line_spacing=1.05)

# bottom callout band
cb_y = diag_top + FIG_H + 0.16
add_rect(s, CONTENT_L, cb_y, CONTENT_W, 0.92, NAVY)
add_rect(s, CONTENT_L, cb_y, 0.14, 0.92, GOLD)
add_text(s, CONTENT_L + 0.35, cb_y + 0.09, CONTENT_W - 0.7, 0.76,
         [[("Dense:", dict(size=16, bold=True, color=GOLD)),
           ("  quality and cost both grow with parameter count — scaling gets expensive fast.", dict(size=16, color=WHITE))],
          [("MoE:", dict(size=16, bold=True, color=GOLD)),
           ("  quality scales with ", dict(size=16, color=WHITE)),
           ("total", dict(size=16, bold=True, italic=True, color=WHITE)),
           (" params, cost only with ", dict(size=16, color=WHITE)),
           ("active", dict(size=16, bold=True, italic=True, color=WHITE)),
           (" params — Qwen3-30B fires just 3B of 30B per token.", dict(size=16, color=WHITE))]],
         space_after=4, line_spacing=1.03)
# tiny source credit — in the gap between the callout band and the footer
add_text(s, SW - 5.8, cb_y + 0.96, 5.2, 0.22,
         [[("Diagram: GShard (Lepikhin et al.), via the Hugging Face MoE guide.",
            dict(size=9, italic=True, color=MUTE))]], align=PP_ALIGN.RIGHT)


# ===========================================================================
# SLIDE 2 — MoE on Edge: the memory-hierarchy challenge
# ===========================================================================
s, top = content_slide("The Challenge: Running MoE on Edge Devices",
                       "At decode, the bottleneck is memory bandwidth — not compute")

# LEFT: memory hierarchy + latency figure
fw, fh = add_figure(s, HIER_FIG, top + 0.35, max_w=6.7, max_h=4.4,
                    center_x=CONTENT_L - 0.05)

# RIGHT: explanation + two challenges
tx0 = CONTENT_L + fw + 0.2
tw = CONTENT_R - tx0

bullets(s, tx0, top + 0.02, tw, 2.0, [
    ([("Decoding is one token at a time.", dict(bold=True, color=NAVY)),
      (" Each token streams its experts from DRAM.", {})], 0),
    ([("Fetching weights dwarfs computing with them", dict(bold=True, color=TEAL)),
      (" — the GPU idles waiting on memory.", {})], 0),
    ([("On edge (phone, L4, Cor3),", dict(bold=True, color=NAVY)),
      (" DRAM is small and slow — both hurt.", {})], 0),
], size=15, space_after=7, line_spacing=1.03)

add_text(s, tx0, top + 1.95, tw, 0.32,
         [[("→ Two ways to make MoE fit and run on edge:",
            dict(size=16, bold=True, color=NAVY))]])

# Challenge 1 card
c1y = top + 2.36
c1h = 0.92
add_rect(s, tx0, c1y, tw, c1h, WHITE, line=RGBColor(0xD5, 0xDE, 0xE7))
add_rect(s, tx0, c1y, 0.12, c1h, GOLD)
add_text(s, tx0 + 0.28, c1y + 0.09, tw - 0.42, c1h - 0.15,
         [[("① Reduce ", dict(size=16, bold=True, color=NAVY)),
           ("total", dict(size=16, bold=True, italic=True, color=TEAL)),
           (" parameters", dict(size=16, bold=True, color=NAVY))],
          [("Shrink the model so it fits in device DRAM.", dict(size=13, color=INK))],
          [("30B = 60 GB, but edge DRAM is only 8–24 GB.", dict(size=12, color=MUTE))]],
         space_after=2, line_spacing=1.0)

# Challenge 2 card
c2y = c1y + c1h + 0.14
c2h = 0.92
add_rect(s, tx0, c2y, tw, c2h, WHITE, line=RGBColor(0xD5, 0xDE, 0xE7))
add_rect(s, tx0, c2y, 0.12, c2h, GOLD)
add_text(s, tx0 + 0.28, c2y + 0.09, tw - 0.42, c2h - 0.15,
         [[("② Reduce ", dict(size=16, bold=True, color=NAVY)),
           ("active", dict(size=16, bold=True, italic=True, color=TEAL)),
           (" parameters", dict(size=16, bold=True, color=NAVY))],
          [("Move fewer bytes into the GPU each token.", dict(size=13, color=INK))],
          [("Fetch/read time scales with bytes → lower latency.", dict(size=12, color=MUTE))]],
         space_after=2, line_spacing=1.0)

add_text(s, tx0, c2y + c2h + 0.06, tw, 0.3,
         [[("This work targets ", dict(size=13, color=INK)),
           ("both", dict(size=13, bold=True, color=NAVY)),
           (".", dict(size=13, color=INK))]])


# ===========================================================================
# Motivation slides (3-6) — the "why per-token, per-channel" argument.
# Each: a headline subtitle, one prominent figure, and a short navy callout
# band carrying the 1-2 key sentences (2-3 sentences total on the page).
# ===========================================================================
def callout(slide, sentences, y=None, h=0.98):
    """Bottom navy band with 1-2 concise sentences (list of run-lists)."""
    if y is None:
        y = CONTENT_BOT - h
    add_rect(slide, CONTENT_L, y, CONTENT_W, h, NAVY)
    add_rect(slide, CONTENT_L, y, 0.14, h, GOLD)
    add_text(slide, CONTENT_L + 0.35, y + 0.10, CONTENT_W - 0.7, h - 0.18,
             sentences, space_after=5, line_spacing=1.05,
             anchor=MSO_ANCHOR.MIDDLE)


# --- Slide 3: a sparse subset of channels suffices -------------------------
s, top = content_slide(
    "A Sparse Subset of Channels Suffices",
    "Within one expert, a few channels carry almost all the output")
cb_y = CONTENT_BOT - 1.02
add_figure(s, FIGS / "fig_sparse_suffices.png", top + 0.10,
           max_w=CONTENT_W, max_h=cb_y - top - 0.28)
callout(s, [
    [("An expert's SwiGLU activations are long-tailed — ", dict(size=17, color=WHITE)),
     ("43% are ≈0", dict(size=17, bold=True, color=GOLD)),
     (", and a thin large-magnitude tail dominates the output.", dict(size=17, color=WHITE))],
    [("Zeroing the smallest 80% of activations barely changes the result: most "
      "channels are idle for any given token.", dict(size=17, color=WHITE))],
], y=cb_y)


# --- Slide 4: the subset is token-specific ---------------------------------
s, top = content_slide(
    "…but the Subset is Token-Specific",
    "Which channels matter is re-decided on every token")
cb_y = CONTENT_BOT - 1.02
add_figure(s, FIGS / "fig_token_specific.png", top + 0.10,
           max_w=CONTENT_W, max_h=cb_y - top - 0.28)
callout(s, [
    [("Consecutive tokens routed to the same expert share only ", dict(size=17, color=WHITE)),
     ("33% of their kept channels", dict(size=17, bold=True, color=GOLD)),
     (" — barely above the 25% random baseline.", dict(size=17, color=WHITE))],
    [("No fixed keep-set works, and even reusing the previous token's mask fails: "
      "the choice must be made per token.", dict(size=17, color=WHITE))],
], y=cb_y)


# --- Slide 5: near-dense accuracy at scale (two figures) -------------------
s, top = content_slide(
    "Acting on It Holds Near-Dense Accuracy",
    "Per-token channel selection on the full dense Qwen3-30B-A3B — no fine-tuning")
cb_y = CONTENT_BOT - 1.02
# two prune-sweep panels side by side, shared height
fmm = Image.open(FIGS / "fig_prune_sweep_mmlu.png")
ar = fmm.width / fmm.height
fig_h = min(3.15, cb_y - top - 0.30)
fig_w = fig_h * ar
gap = 0.45
total = 2 * fig_w + gap
x0 = (SW - total) / 2
fig_y = top + 0.14
add_picture_hw(s, FIGS / "fig_prune_sweep_mmlu.png", x0, fig_y, fig_w, fig_h)
add_picture_hw(s, FIGS / "fig_prune_sweep_hellaswag.png", x0 + fig_w + gap,
               fig_y, fig_w, fig_h)
callout(s, [
    [("Per-token selection stays within noise of dense out to a ", dict(size=17, color=WHITE)),
     ("70% cut", dict(size=17, bold=True, color=GOLD)),
     (" of the active intermediate dimension, on both MMLU and HellaSwag.", dict(size=17, color=WHITE))],
    [("Accuracy bends only past 80% — a large active-parameter cut, essentially "
      "for free.", dict(size=17, color=WHITE))],
], y=cb_y)


# --- Slide 6: online beats offline -----------------------------------------
s, top = content_slide(
    "Online Selection Beats Offline, Decisively",
    "The per-token selector stays near dense where static methods collapse")
cb_y = CONTENT_BOT - 1.02
add_figure(s, FIGS / "fig_offline_vs_online.png", top + 0.10,
           max_w=CONTENT_W, max_h=cb_y - top - 0.28)
callout(s, [
    [("At equal budget, dropping experts or ranking channels ", dict(size=17, color=WHITE)),
     ("offline", dict(size=17, bold=True, color=GOLD)),
     (" falls apart as the cut deepens — toward chance by 7/8.", dict(size=17, color=WHITE))],
    [("Per-token selection holds the dense line; the gap widens to ", dict(size=17, color=WHITE)),
     ("+33 pt", dict(size=17, bold=True, color=GOLD)),
     (" at the tightest budget — the information offline methods cannot see.", dict(size=17, color=WHITE))],
], y=cb_y)


# ===========================================================================
# LEVEL 1 — Offline static channel ranking ("failing experiments")
# Sourced from midpoint_level1.md.  Same navy/gold design system.
# ===========================================================================
CEN = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 8   # generic per-col alignment


def req_card(slide, x, y, w, h, num, title_runs, desc):
    """A white requirement/mechanism card with a gold left spine + number."""
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xD5, 0xDE, 0xE7))
    add_rect(slide, x, y, 0.12, h, GOLD)
    add_text(slide, x + 0.26, y + 0.10, w - 0.42, h - 0.18,
             [[(num + "  ", dict(size=15, bold=True, color=GOLD))] + title_runs,
              [(desc, dict(size=13, color=INK))]],
             space_after=3, line_spacing=1.02)


# --- Section divider -------------------------------------------------------
section_slide("Level 1 — Offline Static Channel Ranking",
              "Can a fixed, precomputed channel order replace the online scorer?")

# --- Slide L1.1: the target ------------------------------------------------
s, top = content_slide(
    "The Target: a Fixed Ranking of Unique Channels",
    "If importance were token-independent, we could rank once and skip the online scorer")
add_text(s, CONTENT_L, top - 0.02, CONTENT_W, 0.9,
         [[("The online method pins ", dict(size=16, color=INK)),
           ("up_proj", dict(size=16, bold=True, color=TEAL)),
           (" at full width — it ", dict(size=16, color=INK)),
           ("is", dict(size=16, italic=True, color=INK)),
           (" the scorer. Offline question: precompute a fixed per-expert channel "
            "order so inference just keeps the top-B, with ", dict(size=16, color=INK)),
           ("no full-width computation", dict(size=16, bold=True, color=NAVY)),
           ("?", dict(size=16, color=INK))]],
         line_spacing=1.05)

# 2x2 grid of the four requirements for a viable offline score
gx, gy = CONTENT_L, top + 0.98
cw = (CONTENT_W - 0.30) / 2
ch, cgap = 1.28, 0.16
reqs = [
    ("①", [("Redundancy-aware", dict(size=15, bold=True, color=NAVY))],
     "Don't double-spend budget on duplicate channels."),
    ("②", [("Nested order", dict(size=15, bold=True, color=NAVY))],
     "Every prefix must be good — the budget varies per token."),
    ("③", [("Cross-expert comparable", dict(size=15, bold=True, color=NAVY))],
     "One global threshold; per-expert quotas emerge on their own."),
    ("④", [("Reweightable by router", dict(size=15, bold=True, color=NAVY))],
     "The only free online signal is the gate weight g_e(x)."),
]
for i, (num, tr, desc) in enumerate(reqs):
    r, c = divmod(i, 2)
    req_card(s, gx + c * (cw + 0.30), gy + r * (ch + cgap), cw, ch, num, tr, desc)

callout(s, [
    [("Payoff if it worked: ", dict(size=16, bold=True, color=GOLD)),
     ("all three matrices narrow, “keep top-B” is a contiguous prefix slice, "
      "and the online cost is ≈ 0.", dict(size=16, color=WHITE))],
], y=gy + 2 * ch + cgap + 0.16, h=0.78)

# --- Slide L1.2: pivoted-Cholesky ------------------------------------------
s, top = content_slide(
    "Pivoted-Cholesky: Greedy, Conditional Selection",
    "Pick the most important channel, then subtract what it already explains")
bullets(s, CONTENT_L, top, CONTENT_W, 1.2, [
    ([("Build a per-expert ", {}), ("coupling matrix", dict(bold=True, color=NAVY)),
      (" (activation covariance ⊙ weight Gram), then run pivoted Cholesky: greedily "
       "select the largest-residual channel and ", {}),
      ("downdate", dict(bold=True, color=TEAL)),
      (" every remaining channel by what the chosen one explains.", {})], 0),
    ([("Duplicates collapse to ≈ 0", dict(bold=True, color=NAVY)),
      (" — a channel already explained is never picked again.", {})], 0),
], size=16, space_after=8, line_spacing=1.05)

# online scoring rule — mono-spaced highlight box
box_y = top + 1.55
add_rect(s, CONTENT_L, box_y, CONTENT_W, 0.86, WHITE, line=RGBColor(0xD5, 0xDE, 0xE7))
add_rect(s, CONTENT_L, box_y, 0.14, 0.86, GOLD)
add_text(s, CONTENT_L + 0.32, box_y + 0.10, CONTENT_W - 0.6, 0.66,
         [[("Online score:   g²(x) · σ", dict(size=19, name="Consolas", bold=True, color=NAVY)),
           ("e,r", dict(size=13, name="Consolas", bold=True, color=NAVY)),
           ("        keep global top-B", dict(size=19, name="Consolas", bold=True, color=NAVY))],
          [("router weight  ×  precomputed marginal gain", dict(size=13, name="Consolas", color=MUTE))]],
         space_after=3, line_spacing=1.0)

bullets(s, CONTENT_L, box_y + 1.02, CONTENT_W, 1.9, [
    ([("Pivot order is ", {}), ("nested and monotone", dict(bold=True, color=NAVY)),
      (" → a single global threshold cuts a clean prefix.", {})], 0),
    ([("Budget-agnostic: stored once (~57 MB), no weight modified.", {})], 0),
    ([("Overhead ≈ ", {}), ("0.016%", dict(bold=True, color=TEAL)),
      (" of expert-FFN MACs.", {})], 0),
], size=16, space_after=8, line_spacing=1.05)

callout(s, [
    [("Every requirement met by construction — a redundancy-aware, nested, "
      "cross-expert-comparable, router-reweightable ranking.", dict(size=17, color=WHITE))],
], h=0.80)

# --- Slide L1.3: results ----------------------------------------------------
s, top = content_slide(
    "Results: Best Offline Ranking, Still Beaten Online",
    "Pivoted-Cholesky tops the offline bracket — yet trails per-token selection by 4–33 pt")
# left: table (hero)
tbl_rows = [
    ["Active cut", "Top-k", "MoSE", "Piv-Chol", "Online"],
    ["Dense", "—", "—", "—", "78.56"],
    ["−50%", "75.2", "69.45", "74.26", "78.54"],
    ["−62.5%", "69.8", "61.00", "70.54", "78.76"],
    ["−75%", "49.4", "43.66", "63.60", "78.28"],
    ["−87.5%", "26.2", "30.32", "44.15", "76.84"],
]
tw = 6.05
_, th = add_table(s, CONTENT_L, top + 0.02, tw, tbl_rows,
                  col_w=[1.30, 1.12, 1.18, 1.30, 1.15], font=14,
                  highlight_cols=(4,), col_align=CEN,
                  row_h=0.40, header_h=0.42,
                  bold_cells={(2, 3), (3, 3), (4, 3), (5, 3)})
add_text(s, CONTENT_L, top + 0.06 + th, tw, 0.3,
         [[("Top-k = drop experts 8→k;  Piv-Chol / MoSE = fixed offline rankings;  "
            "Online = per-token.  HellaSwag acc_norm, no fine-tuning.",
            dict(size=10, italic=True, color=MUTE))]], line_spacing=1.0)
bullets(s, CONTENT_L, top + 0.06 + th + 0.42, tw, 2.4, [
    ([("Pivoted-Cholesky dominates the offline bracket ", dict(bold=True, color=NAVY)),
      ("(+5 to +20 pt over MoSE).", {})], 0),
    ([("But the whole offline family caps out — trails online by ", {}),
      ("4 / 15 / 33 pt", dict(bold=True, color=RED)), (".", {})], 0),
    ([("Cross-expert offline coupling buys nothing (<2% selection change).", {})], 0),
    ([("The headroom is ", {}),
      ("per-token activation information", dict(bold=True, color=NAVY)),
      (" no fixed ranking can capture.", {})], 0),
], size=13, space_after=6, line_spacing=1.02)
# right: the offline-vs-online figure
fim = Image.open(FIGS / "fig_offline_vs_online.png")
ar = fim.width / fim.height
fx = CONTENT_L + tw + 0.30
fw = CONTENT_R - fx
fh = fw / ar
fy = top + 0.10
add_picture_hw(s, FIGS / "fig_offline_vs_online.png", fx, fy, fw, fh)
add_text(s, fx, fy + fh + 0.06, fw, 0.4,
         [[("→ The deployable method must score online.",
            dict(size=15, bold=True, color=NAVY))]], align=PP_ALIGN.CENTER)


# ===========================================================================
# Bonus results
# ===========================================================================
section_slide("Bonus Results",
              "Online selection composes, a threshold beats top-B, and where the cloud win really is")

# --- Slide L1.4: stacks on a compressed base -------------------------------
s, top = content_slide(
    "Online Selection Stacks on a Compressed Base",
    "On a 33%-Nyström + KD-healed base, per-token selection degrades like the dense model")
rows = [
    ["nominal", "online cut", "wikitext ppl ↓", "mmlu ↑", "hellaswag ↑", "gsm8k ↑"],
    ["base", "0%", "10.11", "0.767", "0.799", "0.817"],
    ["50%", "−32.6%", "10.30", "0.763", "0.795", "0.820"],
    ["70%", "−45.7%", "10.96", "0.749", "0.786", "0.787"],
    ["80%", "−52.2%", "11.90", "0.733", "0.767", "0.748"],
]
tw = 9.6
tx = (SW - tw) / 2
_, th = add_table(s, tx, top + 0.25, tw, rows,
                  col_w=[1.35, 1.75, 2.05, 1.45, 1.85, 1.15], font=15,
                  col_align=CEN, row_h=0.46, header_h=0.5)
callout(s, [
    [("Cumulative ", dict(size=17, color=WHITE)),
     ("~68% reduction", dict(size=17, bold=True, color=GOLD)),
     (" of the original expert weights at 80% nominal — for only ", dict(size=17, color=WHITE)),
     ("+1.79 ppl", dict(size=17, bold=True, color=GOLD)),
     (".", dict(size=17, color=WHITE))],
    [("The two axes — offline weight compression + online channel selection — "
      "stack independently.", dict(size=17, color=WHITE))],
], h=1.02)

# --- Slide L1.5: fixed per-layer threshold ---------------------------------
s, top = content_slide(
    "A Fixed Per-Layer Threshold Beats Top-B",
    "A per-layer score threshold beats pooled top-B on every metric at matched budget")
add_text(s, CONTENT_L, top, CONTENT_W, 0.9,
         [[("A threshold ", dict(size=16, color=INK)),
           ("keep iff score ≥ τ", dict(size=16, bold=True, name="Consolas", color=NAVY)),
           (" is ", dict(size=16, color=INK)),
           ("elementwise", dict(size=16, bold=True, color=TEAL)),
           (" — each channel decided independently, no cross-expert sync. Budget "
            "floats per token: hot tokens keep more, quiet tokens fewer.", dict(size=16, color=INK))]],
         line_spacing=1.05)
rows = [
    ["selection rule (80% nominal, matched budget)", "ppl ↓", "mmlu ↑", "hellaswag ↑", "winogrande ↑"],
    ["online top-B (B=1229)", "12.65", "0.779", "0.758", "0.671"],
    ["fixed per-layer threshold", "12.32", "0.785", "0.762", "0.691"],
    ["threshold + FloE predictor", "12.98", "0.771", "0.745", "0.690"],
]
tw = 10.6
tx = (SW - tw) / 2
_, th = add_table(s, tx, top + 1.05, tw, rows,
                  col_w=[4.20, 1.45, 1.55, 1.90, 1.50], font=15,
                  highlight_rows=(2,), col_align=CEN, row_h=0.48, header_h=0.62)
callout(s, [
    [("Threshold wins by ", dict(size=17, color=WHITE)),
     ("reallocating budget across tokens", dict(size=17, bold=True, color=GOLD)),
     (" — variable spend > fixed spend.", dict(size=17, color=WHITE))],
    [("It does not compose with predict-ahead: a stale signal moves both the count "
      "and the choice.", dict(size=17, color=WHITE))],
], h=1.02)

# --- Slide L1.6: cloud resident decode -------------------------------------
s, top = content_slide(
    "Cloud Decode: Dynamic Selection Is the Wrong Lever",
    "On resident GPUs, dynamic selection hurts throughput — tensor parallelism is the win")
# left: throughput table
rows = [
    ["Parallelism", "Dense", "Dynamic", ""],
    ["FSDP2 (dp)", "base", "0.83–0.88×", "↓"],
    ["TP = 2", "~2.2× base", "0.77× dense", "↓"],
    ["TP = 4", "~2.4× base", "0.65× dense", "↓"],
]
tw = 5.9
_, th = add_table(s, CONTENT_L, top + 0.05, tw, rows,
                  col_w=[1.70, 1.70, 1.85, 0.65], font=14,
                  col_align=CEN, row_h=0.5, header_h=0.5)
add_text(s, CONTENT_L, top + 0.10 + th, tw, 0.3,
         [[("Throughput vs. the dense baseline for the same layout (tok/s).",
            dict(size=10, italic=True, color=MUTE))]])
# right: why it hurts
add_text(s, CONTENT_L + tw + 0.35, top - 0.02, CONTENT_R - (CONTENT_L + tw + 0.35), 0.35,
         [[("Why it hurts instead of helps:", dict(size=16, bold=True, color=NAVY))]])
bullets(s, CONTENT_L + tw + 0.35, top + 0.38,
        CONTENT_R - (CONTENT_L + tw + 0.35), 3.0, [
    ([("FSDP2: ", dict(bold=True, color=TEAL)),
      ("the all-gather of each block's sharded weights (~83% of GPU-busy time) is "
       "independent of per-token budget.", {})], 0),
    ([("TP: ", dict(bold=True, color=TEAL)),
      ("the weight read is already split tp-ways to a tiny slice (18.9 MB/rank at "
       "tp=4) — nothing left worth cutting.", {})], 0),
    ([("Fixed overhead added: ", dict(bold=True, color=RED)),
      ("~0.17 ms/token top-k + gather/scatter, plus ~0.1 ms all-reduce over H under TP.", {})], 0),
], size=13.5, space_after=8, line_spacing=1.04)
callout(s, [
    [("Cloud throughput levers are ", dict(size=16, color=WHITE)),
     ("collectives + batch/parallelism layout", dict(size=16, bold=True, color=GOLD)),
     (", not activated-parameter count.", dict(size=16, color=WHITE))],
    [("The dynamic method's cloud payoff is ", dict(size=16, color=WHITE)),
     ("activation memory, not tokens/s", dict(size=16, bold=True, color=GOLD)),
     (" — it targets edge, not cloud.", dict(size=16, color=WHITE))],
], h=1.02)


# ===========================================================================
prs.save(str(OUT))
print(f"\nSaved {OUT} with {len(prs.slides)} slides")
